"""
AI Study Assistant - Flask Application
Supports SQLite (local) and PostgreSQL (production via DATABASE_URL)
"""

from flask import Flask, render_template, request, jsonify, session
from werkzeug.utils import secure_filename
from flask_cors import CORS
import os
import json
import re
import uuid
import sqlite3
from datetime import datetime
from typing import Dict

# ── Optional dependencies ──────────────────────────────────────────────────
try:
    from groq import Groq
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False
    print("⚠️  groq not installed. Run: pip install groq")

try:
    import PyPDF2 as pypdf
    PYPDF_AVAILABLE = True
except ImportError:
    PYPDF_AVAILABLE = False

try:
    from docx import Document as DocxDocument
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

try:
    from pptx import Presentation as PptxPresentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed. Run: pip install python-pptx")

# ── App setup ──────────────────────────────────────────────────────────────
app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'dev-secret-change-in-production')
CORS(app)

# ── Config from environment ────────────────────────────────────────────────
GROQ_API_KEY     = os.environ.get('GROQ_API_KEY', '')
DATABASE_URL     = os.environ.get('DATABASE_URL', '')
DATABASE_FILE    = os.environ.get('DATABASE_FILE', 'study_data.db')
# Use /tmp on cloud (ephemeral), local folder otherwise
IS_PRODUCTION    = bool(os.environ.get('RENDER') or os.environ.get('RAILWAY_ENVIRONMENT'))
UPLOAD_FOLDER    = '/tmp/uploads' if IS_PRODUCTION else 'uploads'
ALLOWED_EXTENSIONS = {'pptx', 'txt', 'docx', 'doc'}
MAX_FILE_SIZE    = 400 * 1024 * 1024  # 400 MB
MAX_CONTENT_PER_DOC = 20000

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

# Determine database dialect
USE_POSTGRES = DATABASE_URL.startswith(('postgresql', 'postgres'))
PH = '%s' if USE_POSTGRES else '?'   # parameter placeholder

# ── Groq client ───────────────────────────────────────────────────────────
GROQ_MODEL = 'llama-3.3-70b-versatile'
groq_client = Groq(api_key=GROQ_API_KEY) if (GROQ_AVAILABLE and GROQ_API_KEY) else None

# In-memory document store (cache loaded from DB on startup / upload)
documents: Dict = {}

# ── Database helpers ───────────────────────────────────────────────────────
def get_conn():
    """Return a DB connection (PostgreSQL or SQLite)."""
    if USE_POSTGRES:
        import psycopg2
        url = DATABASE_URL.replace('postgres://', 'postgresql://', 1)
        return psycopg2.connect(url, sslmode='require')
    return sqlite3.connect(DATABASE_FILE)


def init_db():
    """Create all tables if they don't exist."""
    conn = get_conn()
    cur = conn.cursor()

    cur.execute('''
        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            content TEXT NOT NULL,
            upload_time TEXT NOT NULL,
            file_size INTEGER,
            content_preview TEXT
        )
    ''')

    cur.execute('''
        CREATE TABLE IF NOT EXISTS questions (
            id TEXT PRIMARY KEY,
            doc_id TEXT NOT NULL,
            question_type TEXT,
            questions_json TEXT NOT NULL,
            created_time TEXT NOT NULL
        )
    ''')

    cur.execute('''
        CREATE TABLE IF NOT EXISTS summaries (
            id TEXT PRIMARY KEY,
            doc_id TEXT NOT NULL,
            summary_json TEXT NOT NULL,
            created_time TEXT NOT NULL
        )
    ''')

    cur.execute('''
        CREATE TABLE IF NOT EXISTS chat_messages (
            id TEXT PRIMARY KEY,
            session_id TEXT NOT NULL,
            doc_ids TEXT,
            role TEXT NOT NULL,
            content TEXT NOT NULL,
            created_time TEXT NOT NULL
        )
    ''')

    conn.commit()
    conn.close()


def _upsert_document(cur, doc_id, name, content, file_size, now, preview):
    if USE_POSTGRES:
        cur.execute(f'''
            INSERT INTO documents (id, name, content, upload_time, file_size, content_preview)
            VALUES ({PH},{PH},{PH},{PH},{PH},{PH})
            ON CONFLICT (id) DO UPDATE SET
                name=EXCLUDED.name, content=EXCLUDED.content,
                upload_time=EXCLUDED.upload_time, file_size=EXCLUDED.file_size,
                content_preview=EXCLUDED.content_preview
        ''', (doc_id, name, content, now, file_size, preview))
    else:
        cur.execute('''
            INSERT OR REPLACE INTO documents (id, name, content, upload_time, file_size, content_preview)
            VALUES (?,?,?,?,?,?)
        ''', (doc_id, name, content, now, file_size, preview))


def save_document_to_db(doc_id, name, content, file_size):
    conn = get_conn()
    cur = conn.cursor()
    _upsert_document(cur, doc_id, name, content, file_size,
                     datetime.now().isoformat(), content[:200])
    conn.commit()
    conn.close()


def save_questions_to_db(doc_id, questions):
    conn = get_conn()
    cur = conn.cursor()
    q_id = str(uuid.uuid4())
    cur.execute(f'''
        INSERT INTO questions (id, doc_id, question_type, questions_json, created_time)
        VALUES ({PH},{PH},{PH},{PH},{PH})
    ''', (q_id, doc_id, 'all', json.dumps(questions), datetime.now().isoformat()))
    conn.commit()
    conn.close()
    return q_id


def save_summary_to_db(doc_id, summary):
    conn = get_conn()
    cur = conn.cursor()
    s_id = str(uuid.uuid4())
    cur.execute(f'''
        INSERT INTO summaries (id, doc_id, summary_json, created_time)
        VALUES ({PH},{PH},{PH},{PH})
    ''', (s_id, doc_id, json.dumps(summary), datetime.now().isoformat()))
    conn.commit()
    conn.close()
    return s_id


def save_chat_message(session_id, role, content, doc_ids=None):
    conn = get_conn()
    cur = conn.cursor()
    m_id = str(uuid.uuid4())
    cur.execute(f'''
        INSERT INTO chat_messages (id, session_id, doc_ids, role, content, created_time)
        VALUES ({PH},{PH},{PH},{PH},{PH},{PH})
    ''', (m_id, session_id, json.dumps(doc_ids or []), role, content, datetime.now().isoformat()))
    conn.commit()
    conn.close()


def get_chat_history(session_id, limit=50):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute(f'''
        SELECT role, content, created_time
        FROM chat_messages
        WHERE session_id = {PH}
        ORDER BY created_time DESC
        LIMIT {PH}
    ''', (session_id, limit))
    rows = cur.fetchall()
    conn.close()
    return [{'role': r[0], 'content': r[1], 'created_time': r[2]} for r in reversed(rows)]


def get_documents_from_db():
    conn = get_conn()
    cur = conn.cursor()
    cur.execute('SELECT id, name, upload_time, file_size FROM documents ORDER BY upload_time DESC')
    rows = cur.fetchall()
    conn.close()
    return [{'id': r[0], 'name': r[1], 'upload_time': r[2], 'size': r[3]} for r in rows]


def get_document_from_db(doc_id):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute(f'SELECT id, name, content, upload_time, file_size FROM documents WHERE id = {PH}', (doc_id,))
    row = cur.fetchone()
    conn.close()
    if row:
        return {'id': row[0], 'name': row[1], 'content': row[2], 'upload_time': row[3], 'file_size': row[4]}
    return None


def get_questions_for_doc(doc_id):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute(f'''
        SELECT id, questions_json, created_time FROM questions
        WHERE doc_id = {PH} ORDER BY created_time DESC LIMIT 1
    ''', (doc_id,))
    row = cur.fetchone()
    conn.close()
    if row:
        return {'id': row[0], 'questions': json.loads(row[1]), 'created_time': row[2]}
    return None


def get_summary_for_doc(doc_id):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute(f'''
        SELECT id, summary_json, created_time FROM summaries
        WHERE doc_id = {PH} ORDER BY created_time DESC LIMIT 1
    ''', (doc_id,))
    row = cur.fetchone()
    conn.close()
    if row:
        return {'id': row[0], 'summary': json.loads(row[1]), 'created_time': row[2]}
    return None


def delete_document_from_db(doc_id):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute(f'DELETE FROM questions WHERE doc_id = {PH}', (doc_id,))
    cur.execute(f'DELETE FROM summaries WHERE doc_id = {PH}', (doc_id,))
    cur.execute(f'DELETE FROM documents WHERE id = {PH}', (doc_id,))
    conn.commit()
    conn.close()


def load_documents_from_db():
    """Populate in-memory cache from DB (called at startup)."""
    global documents
    conn = get_conn()
    cur = conn.cursor()
    cur.execute('SELECT id, name, content, upload_time, file_size FROM documents')
    for row in cur.fetchall():
        documents[row[0]] = {
            'name': row[1],
            'content': row[2],
            'upload_time': row[3],
            'file_size': row[4],
            'size': len(row[2]),
            'filepath': None,
        }
    conn.close()
    print(f"✓ Loaded {len(documents)} document(s) from database")

# ── File helpers ───────────────────────────────────────────────────────────
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def extract_text_from_file(filepath):
    ext = filepath.rsplit('.', 1)[1].lower()
    text = ''

    if ext == 'pptx':
        if not PPTX_AVAILABLE:
            raise Exception('python-pptx not installed. Run: pip install python-pptx')
        prs = PptxPresentation(filepath)
        for i, slide in enumerate(prs.slides, 1):
            text += f'\n--- Slide {i} ---\n'
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text += shape.text.strip() + '\n'

    elif ext == 'txt':
        try:
            with open(filepath, encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            with open(filepath, encoding='latin-1') as f:
                text = f.read()

    elif ext in ('docx', 'doc'):
        if not PYTHON_DOCX_AVAILABLE:
            raise Exception('python-docx not installed')
        doc = DocxDocument(filepath)
        text = '\n'.join(p.text for p in doc.paragraphs if p.text.strip())

    text = re.sub(r'\n{3,}', '\n\n', text).strip()
    if not text:
        raise Exception('File is empty or contains no extractable text')
    return text


def clean_text(text, max_length=MAX_CONTENT_PER_DOC):
    text = ' '.join(text.split())
    return text[:max_length]


def parse_json_response(response_text):
    try:
        start = response_text.find('{')
        end = response_text.rfind('}') + 1
        if start != -1 and end > start:
            return json.loads(response_text[start:end])
        return json.loads(response_text)
    except json.JSONDecodeError as e:
        return {'open_ended': [], 'multiple_choice': [], 'true_false': [], 'fill_in_blank': [],
                'error': f'Failed to parse JSON: {e}'}

# ── AI helpers ─────────────────────────────────────────────────────────────
def generate_questions_from_text(content):
    if not groq_client:
        return {'error': 'Groq API not configured. Set GROQ_API_KEY environment variable.'}

    content = clean_text(content)
    if len(content) < 100:
        return {'error': 'Document content too short'}

    prompt = f"""Based on the following document, generate questions in 4 categories.

DOCUMENT:
---
{content}
---

Generate:
- 15 open-ended questions
- 15 multiple choice questions
- 10 true/false questions
- 10 fill-in-the-blank questions

Return ONLY valid JSON:

{{
  "open_ended": [{{"question": "?", "answer": "Answer", "keywords": ["k1"]}}],
  "multiple_choice": [{{"question": "?", "options": ["A","B*","C","D"], "answer": "1", "explanation": "Why", "keywords": ["k1"]}}],
  "true_false": [{{"question": "Statement.", "correct": true, "explanation": "Why", "keywords": ["k1"]}}],
  "fill_in_blank": [{{"question": "The ___ is ___", "answer": "word", "explanation": "Context", "keywords": ["k1"]}}]
}}"""

    try:
        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{'role': 'user', 'content': prompt}]
        )
        return parse_json_response(response.choices[0].message.content)
    except Exception as e:
        return {'error': f'AI generation failed: {e}'}

# ── Routes ─────────────────────────────────────────────────────────────────
@app.route('/')
def index():
    return render_template('index.html')


@app.route('/favicon.ico')
def favicon():
    return jsonify({'status': 'ok'}), 200


@app.errorhandler(413)
def file_too_large(e):
    mb = MAX_FILE_SIZE // (1024 * 1024)
    return jsonify({'error': f'File too large. Maximum allowed size is {mb} MB.'}), 413


@app.route('/api/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file provided'}), 400
        file = request.files['file']
        if not file.filename:
            return jsonify({'error': 'No file selected'}), 400
        if not allowed_file(file.filename):
            return jsonify({'error': f'Unsupported file type. Allowed: {", ".join(ALLOWED_EXTENSIONS)}'}), 400

        filename = secure_filename(file.filename)
        unique_name = f'{uuid.uuid4()}_{filename}'
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_name)
        file.save(filepath)

        try:
            content = extract_text_from_file(filepath)
        except Exception as e:
            try:
                os.remove(filepath)
            except Exception:
                pass
            return jsonify({'error': f'Could not extract text: {e}'}), 400

        if len(content) < 20:
            os.remove(filepath)
            return jsonify({'error': 'File appears empty or unreadable'}), 400

        doc_id = str(uuid.uuid4())
        file_size = os.path.getsize(filepath)
        documents[doc_id] = {
            'name': file.filename,
            'filepath': filepath,
            'content': content,
            'upload_time': datetime.now().isoformat(),
            'size': len(content),
            'file_size': file_size,
        }
        save_document_to_db(doc_id, file.filename, content, file_size)

        print(f'✓ Uploaded: {file.filename} (ID: {doc_id})')
        return jsonify({
            'success': True,
            'doc_id': doc_id,
            'filename': file.filename,
            'content_preview': content[:150],
            'content_length': len(content),
            'message': f'Successfully uploaded {file.filename}',
        }), 200

    except Exception as e:
        return jsonify({'error': f'Upload failed: {e}'}), 500


@app.route('/api/documents', methods=['GET'])
def get_documents():
    docs = get_documents_from_db()
    if not docs:
        docs = [{'id': did, 'name': d['name'], 'upload_time': d['upload_time'], 'size': d['size']}
                for did, d in documents.items()]
    return jsonify({'documents': docs}), 200


@app.route('/api/document-content/<doc_id>', methods=['GET'])
def get_document_content(doc_id):
    doc = get_document_from_db(doc_id) or (
        {'id': doc_id, **documents[doc_id]} if doc_id in documents else None
    )
    if not doc:
        return jsonify({'error': 'Document not found'}), 404
    return jsonify({'success': True, **doc}), 200


@app.route('/api/generate-questions', methods=['POST'])
def generate_questions():
    try:
        data = request.json or {}
        doc_id = data.get('doc_id')

        # Try memory first, then DB
        if doc_id not in documents:
            db_doc = get_document_from_db(doc_id)
            if db_doc:
                documents[doc_id] = db_doc
                documents[doc_id]['size'] = len(db_doc['content'])
            else:
                return jsonify({'error': 'Document not found'}), 404

        content = documents[doc_id]['content']
        questions = generate_questions_from_text(content)

        if 'error' not in questions:
            save_questions_to_db(doc_id, questions)

        return jsonify({
            'success': True,
            'doc_name': documents[doc_id]['name'],
            'questions': questions,
        }), 200

    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/chat', methods=['POST'])
def chat():
    try:
        if not groq_client:
            return jsonify({'error': 'Groq API not configured. Set GROQ_API_KEY environment variable.'}), 503

        data = request.json or {}
        message = (data.get('message') or '').strip()
        doc_id = data.get('doc_id')
        use_all_docs = data.get('use_all_docs', True)

        if not message:
            return jsonify({'error': 'Message cannot be empty'}), 400
        if len(message) > 2000:
            return jsonify({'error': 'Message too long (max 2000 chars)'}), 400
        if not documents:
            return jsonify({'error': 'No documents uploaded. Please upload a document first.'}), 400

        # Session ID for chat history
        if 'chat_session_id' not in session:
            session['chat_session_id'] = str(uuid.uuid4())
        session_id = session['chat_session_id']

        # Build document context
        context = ''
        doc_names = []
        used_doc_ids = []

        if doc_id and doc_id in documents:
            context = documents[doc_id]['content'][:3000]
            doc_names.append(documents[doc_id]['name'])
            used_doc_ids.append(doc_id)
        elif use_all_docs:
            for did, doc in documents.items():
                doc_names.append(doc['name'])
                used_doc_ids.append(did)
                context += f"\n[From {doc['name']}]\n" + doc['content'][:1500]
        else:
            return jsonify({'error': 'No document selected'}), 400

        # Retrieve recent history for context
        history = get_chat_history(session_id, limit=10)
        history_text = '\n'.join(
            f"{'User' if m['role']=='user' else 'Assistant'}: {m['content']}"
            for m in history[-6:]
        ) if history else ''

        system_prompt = f"""You are an educational AI assistant helping students learn from their documents.

Rules:
1. Only answer based on the source documents below.
2. If the answer is not in the sources, say so clearly.
3. Cite the document you are referencing.
4. Be concise (2-3 paragraphs max).

SOURCES: {', '.join(doc_names)}

CONTENT:
---
{context}
---
{f"RECENT CONVERSATION:{chr(10)}{history_text}" if history_text else ''}"""

        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[
                {'role': 'system', 'content': system_prompt},
                {'role': 'user', 'content': message},
            ]
        )
        bot_reply = response.choices[0].message.content

        # Persist both turns
        save_chat_message(session_id, 'user', message, used_doc_ids)
        save_chat_message(session_id, 'assistant', bot_reply, used_doc_ids)

        return jsonify({
            'success': True,
            'response': bot_reply,
            'sources_used': doc_names,
            'session_id': session_id,
        }), 200

    except Exception as e:
        return jsonify({'error': f'Chat failed: {e}'}), 500


@app.route('/api/chat-history', methods=['GET'])
def chat_history_endpoint():
    """Return past chat messages for the current session."""
    if 'chat_session_id' not in session:
        return jsonify({'history': []}), 200
    history = get_chat_history(session['chat_session_id'], limit=100)
    return jsonify({'history': history, 'session_id': session['chat_session_id']}), 200


@app.route('/api/summarize', methods=['POST'])
def summarize():
    try:
        if not groq_client:
            return jsonify({'error': 'Groq API not configured. Set GROQ_API_KEY environment variable.'}), 503

        data = request.json or {}
        doc_id = data.get('doc_id')

        if not documents:
            return jsonify({'error': 'No documents uploaded'}), 400

        if doc_id:
            if doc_id not in documents:
                db_doc = get_document_from_db(doc_id)
                if db_doc:
                    documents[doc_id] = db_doc
                    documents[doc_id]['size'] = len(db_doc['content'])
                else:
                    return jsonify({'error': 'Document not found'}), 404
            content = documents[doc_id]['content'][:5000]
            doc_name = documents[doc_id]['name']
        else:
            content = ''
            names = []
            for did, doc in documents.items():
                names.append(doc['name'])
                content += f"\n[From {doc['name']}]\n" + doc['content'][:2000]
            doc_name = f'All Documents ({len(names)})'

        prompt = f"""Summarize the following document comprehensively.

DOCUMENT:
---
{content}
---

Return ONLY valid JSON:
{{
  "title": "Brief document title",
  "overview": "2-3 paragraph overview",
  "main_sections": [
    {{"section_name": "...", "description": "...", "key_concepts": ["..."]}}
  ],
  "key_points": ["point1", "point2", "point3", "point4", "point5"],
  "keyword_definitions": [
    {{"term": "...", "definition": "...", "context": "..."}}
  ],
  "conclusion": "Main takeaways"
}}"""

        response = groq_client.chat.completions.create(
            model=GROQ_MODEL,
            messages=[{'role': 'user', 'content': prompt}]
        )
        summary_data = parse_json_response(response.choices[0].message.content)

        if 'error' not in summary_data and doc_id:
            save_summary_to_db(doc_id, summary_data)

        return jsonify({'success': True, 'document': doc_name, **summary_data}), 200

    except Exception as e:
        return jsonify({'error': f'Summary failed: {e}'}), 500


@app.route('/api/document/<doc_id>', methods=['GET'])
def get_document(doc_id):
    if doc_id not in documents:
        return jsonify({'error': 'Document not found'}), 404
    doc = documents[doc_id]
    return jsonify({
        'id': doc_id, 'name': doc['name'],
        'upload_time': doc['upload_time'], 'size': doc['size'],
        'preview': doc['content'][:500],
    }), 200


@app.route('/api/delete-document/<doc_id>', methods=['DELETE'])
def delete_document(doc_id):
    try:
        if doc_id not in documents:
            return jsonify({'error': 'Document not found'}), 404
        doc = documents[doc_id]
        if doc.get('filepath') and os.path.exists(doc['filepath']):
            os.remove(doc['filepath'])
        del documents[doc_id]
        delete_document_from_db(doc_id)
        return jsonify({'success': True, 'message': 'Document deleted'}), 200
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/health', methods=['GET'])
def health():
    return jsonify({
        'status': 'online',
        'version': '2.0',
        'documents_loaded': len(documents),
        'database': 'postgresql' if USE_POSTGRES else 'sqlite',
        'groq_api': 'configured' if groq_client else 'not configured',
        'groq_package_installed': GROQ_AVAILABLE,
        'groq_key_set': bool(GROQ_API_KEY),
        'features': {
            'pptx_support': PPTX_AVAILABLE,
            'docx_support': PYTHON_DOCX_AVAILABLE,
            'ai_enabled': GROQ_AVAILABLE and bool(GROQ_API_KEY),
        },
        'timestamp': datetime.now().isoformat(),
    }), 200


@app.route('/api/debug', methods=['GET'])
def debug():
    """Diagnostic endpoint — shows what is configured without exposing secrets."""
    key = GROQ_API_KEY
    masked = (key[:4] + '...' + key[-4:]) if len(key) >= 8 else ('(empty)' if not key else '(too short)')
    return jsonify({
        'GROQ_AVAILABLE': GROQ_AVAILABLE,
        'GROQ_API_KEY': masked,
        'groq_client_created': groq_client is not None,
        'DATABASE_URL_set': bool(DATABASE_URL),
        'USE_POSTGRES': USE_POSTGRES,
        'UPLOAD_FOLDER': UPLOAD_FOLDER,
        'IS_PRODUCTION': IS_PRODUCTION,
        'documents_in_memory': len(documents),
    }), 200


@app.route('/api/stats', methods=['GET'])
def stats():
    total_size = sum(d['size'] for d in documents.values())
    return jsonify({
        'total_documents': len(documents),
        'total_content_size': total_size,
        'avg_doc_size': total_size // len(documents) if documents else 0,
        'documents': [{'id': did, 'name': d['name'], 'size': d['size']} for did, d in documents.items()],
    }), 200


# ── Startup ────────────────────────────────────────────────────────────────
def startup():
    print('\n' + '='*50)
    print('   AI Study Assistant')
    print('='*50)
    print(f'  Database : {"PostgreSQL" if USE_POSTGRES else f"SQLite ({DATABASE_FILE})"}')
    print(f'  Groq     : {"Enabled (" + GROQ_MODEL + ")" if groq_client else "Not configured — set GROQ_API_KEY"}')
    print(f'  PPTX     : {"Enabled" if PPTX_AVAILABLE else "Not installed"}')
    print(f'  DOCX/DOC : {"Enabled" if PYTHON_DOCX_AVAILABLE else "Not installed"}')

    init_db()
    print('✓ Database initialized')

    load_documents_from_db()
    print('='*50 + '\n')


startup()


if __name__ == '__main__':
    print('  Open http://127.0.0.1:5000')
    app.run(debug=True, port=5000)
